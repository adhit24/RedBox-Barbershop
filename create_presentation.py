"""
RedBox Barbershop — Generator Presentasi Alur Reservasi
Menghasilkan file: RedBox_Reservasi_Flow.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ──────────────────────────────────────────────────────────────
RED       = RGBColor(0xC0, 0x00, 0x00)   # RedBox merah utama
DARK_RED  = RGBColor(0x8B, 0x00, 0x00)   # Merah gelap aksen
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)   # Hitam latar
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)   # Abu gelap slide biasa
MID_GRAY  = RGBColor(0x44, 0x44, 0x44)   # Kotak konten
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)   # Latar terang
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW    = RGBColor(0xFF, 0xC0, 0x00)   # Aksen highlight
GREEN     = RGBColor(0x00, 0xB0, 0x50)
ORANGE    = RGBColor(0xFF, 0x7B, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # blank layout

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_textbox_multiline(slide, lines, x, y, w, h,
                           font_size=14, bold=False, color=WHITE,
                           align=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def slide_bg(slide, color=DARK_GRAY):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color)
    return bg

def top_bar(slide, title, subtitle=None, bar_color=RED):
    add_rect(slide, 0, 0, 13.33, 1.4, fill_rgb=bar_color)
    add_rect(slide, 0, 1.35, 13.33, 0.08, fill_rgb=YELLOW)
    add_text(slide, title, 0.5, 0.1, 12, 0.7,
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.75, 12, 0.55,
                 font_size=16, color=RGBColor(0xFF, 0xDD, 0xDD), align=PP_ALIGN.LEFT)

def section_box(slide, x, y, w, h, title, lines, title_color=RED, bg=MID_GRAY,
                font_size=13, title_size=15):
    add_rect(slide, x, y, w, h, fill_rgb=bg,
             line_rgb=RGBColor(0x60, 0x60, 0x60), line_width=0.5)
    # Title bar inside box
    add_rect(slide, x, y, w, 0.38, fill_rgb=title_color)
    add_text(slide, title, x + 0.12, y + 0.04, w - 0.2, 0.32,
             font_size=title_size, bold=True, color=WHITE)
    # Content lines
    content_y = y + 0.44
    line_h = (h - 0.5) / max(len(lines), 1)
    line_h = min(line_h, 0.38)
    for i, line in enumerate(lines):
        add_text(slide, line, x + 0.18, content_y + i * line_h,
                 w - 0.25, line_h + 0.05, font_size=font_size, color=WHITE)

def arrow_right(slide, x, y, size=0.3, color=YELLOW):
    """Draw simple → arrow as text"""
    add_text(slide, "→", x, y, 0.4, 0.4,
             font_size=int(size * 72), bold=True, color=color, align=PP_ALIGN.CENTER)

def circle_step(slide, num, x, y, size=0.6, text_size=22, bg=RED):
    r = size / 2
    add_rect(slide, x, y, size, size, fill_rgb=bg)
    add_text(slide, str(num), x, y, size, size,
             font_size=text_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 1 — COVER
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, BLACK)

# Background accent blocks
add_rect(slide, 0, 0, 0.6, 7.5, fill_rgb=RED)
add_rect(slide, 0.6, 0, 12.73, 7.5, fill_rgb=BLACK)
add_rect(slide, 0.6, 3.2, 12.73, 0.06, fill_rgb=RED)

# Logo placeholder text
add_rect(slide, 1.2, 0.6, 2.5, 1.0, fill_rgb=RED)
add_text(slide, "RED BOX", 1.2, 0.6, 2.5, 1.0,
         font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 1.2, 1.55, 2.5, 0.35, fill_rgb=DARK_RED)
add_text(slide, "BARBERSHOP", 1.2, 1.55, 2.5, 0.35,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Main title
add_text(slide, "SISTEM RESERVASI", 1.1, 2.4, 11, 1.1,
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "RedBox Barbershop", 1.1, 3.4, 11, 0.6,
         font_size=26, bold=False, color=RED, align=PP_ALIGN.LEFT)

add_rect(slide, 1.1, 4.1, 6, 0.05, fill_rgb=RGBColor(0x60, 0x60, 0x60))

add_text(slide, "Alur Booking Customer  ·  Tugas Admin per Cabang  ·  Flow Teknis Reservasi",
         1.1, 4.3, 11.5, 0.5,
         font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)

add_text(slide, "2026", 1.1, 5.1, 3, 0.4, font_size=16,
         color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.LEFT)

# =============================================================================
# SLIDE 2 — AGENDA
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Agenda", "Topik yang akan dibahas dalam presentasi ini")

add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

items = [
    ("01", "Alur Booking Customer",
     "Langkah-langkah customer dari website hingga booking terkonfirmasi"),
    ("02", "Notifikasi WhatsApp Otomatis",
     "WA konfirmasi, reminder H-1, dan reminder 1 jam sebelum jadwal"),
    ("03", "Tugas Admin per Cabang",
     "Dashboard, notifikasi masuk, update status, dan koordinasi kapster"),
    ("04", "Flow Teknis Reservasi",
     "Arsitektur sistem: API, database, Moka POS, dan Fonnte gateway"),
    ("05", "Sistem Reminder & Cron",
     "Cron jobs otomatis: H-1 reminder, remind-soon, expire stale bills"),
    ("06", "Status & Lifecycle Booking",
     "Status booking dari pending hingga done/cancel dan alur perubahannya"),
]

cols = [(0.5, 6.4), (6.9, 6.0)]
rows_per_col = 3

for i, (num, title, desc) in enumerate(items):
    col = i // rows_per_col
    row = i % rows_per_col
    cx, cw = cols[col]
    cy = 1.65 + row * 1.85

    add_rect(slide, cx, cy, cw, 1.6,
             fill_rgb=MID_GRAY,
             line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.5)
    add_rect(slide, cx, cy, 0.7, 1.6, fill_rgb=RED)
    add_text(slide, num, cx, cy, 0.7, 1.6,
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, cx + 0.78, cy + 0.15, cw - 0.85, 0.42,
             font_size=16, bold=True, color=WHITE)
    add_text(slide, desc, cx + 0.78, cy + 0.6, cw - 0.85, 0.85,
             font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

# =============================================================================
# SLIDE 3 — ALUR BOOKING CUSTOMER (VISUAL FLOW)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Alur Booking Customer", "Langkah dari buka website hingga booking terkonfirmasi")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

steps = [
    ("1", "Buka Website", "redboxbarbershop.com\n/booking.html"),
    ("2", "Pilih Layanan", "26 layanan tersedia\n(Haircut, Fade, Spa, dll)"),
    ("3", "Pilih Kapster", "Daftar kapster aktif\nper cabang"),
    ("4", "Pilih Tanggal\n& Jam", "Slot real-time\n(tidak bisa double)"),
    ("5", "Isi Data Diri", "Nama lengkap\n& Nomor WA"),
    ("6", "Konfirmasi\nBooking", "Klik Submit\n→ Terkonfirmasi"),
]

box_w = 1.85
gap   = 0.22
total = len(steps) * box_w + (len(steps) - 1) * gap
start_x = (13.33 - total) / 2

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    y = 1.65

    # Card background
    add_rect(slide, x, y, box_w, 3.8,
             fill_rgb=MID_GRAY,
             line_rgb=RGBColor(0x66, 0x66, 0x66), line_width=0.5)

    # Number badge
    add_rect(slide, x + box_w/2 - 0.32, y + 0.18, 0.64, 0.64, fill_rgb=RED)
    add_text(slide, num, x + box_w/2 - 0.32, y + 0.18, 0.64, 0.64,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, title, x + 0.1, y + 0.98, box_w - 0.2, 0.75,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 0.1, y + 1.75, box_w - 0.2, 1.8,
             font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(steps) - 1:
        ax = x + box_w + 0.02
        add_text(slide, "▶", ax, y + 1.65, gap + 0.1, 0.4,
                 font_size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Result box at bottom
add_rect(slide, 1.0, 5.75, 11.33, 1.5, fill_rgb=RGBColor(0x00, 0x60, 0x00),
         line_rgb=GREEN, line_width=1)
add_text(slide, "✓  BOOKING CONFIRMED!", 1.2, 5.88, 11, 0.5,
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Halaman konfirmasi muncul · WA konfirmasi otomatis terkirim · Data tersimpan di sistem",
         1.2, 6.38, 11, 0.4,
         font_size=13, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 4 — NOTIFIKASI WA CUSTOMER
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Notifikasi WhatsApp Otomatis", "3 titik pengiriman WA sepanjang lifecycle booking customer")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

notifs = [
    (GREEN,  "📲  Notif 1 — KONFIRMASI BOOKING",
     "Dikirim: Saat booking berhasil dikonfirmasi (real-time)",
     [
         "• Nama customer + sapaan personal",
         "• Detail layanan, harga, dan durasi",
         "• Tanggal & jam booking (format Indonesia)",
         "• Nama kapster yang dipilih",
         "• Nama cabang RedBox",
         '• Closing: "Ada yang mau ditanyain? aku siap bantu kapan aja! 💬"',
     ]),
    (ORANGE, "🔔  Notif 2 — REMINDER H-1",
     "Dikirim: Setiap hari jam 10:00 WIB (cron harian)",
     [
         "• Pengingat jadwal besok",
         "• Detail tanggal, jam, layanan, kapster",
         "• Nama cabang tujuan",
         "• Pesan: 'Dateng tepat waktu ya kak'",
         "• Berlaku untuk SEMUA cabang otomatis",
         "• Tidak dikirim jika status booking cancelled",
     ]),
    (YELLOW, "⏰  Notif 3 — REMIND-SOON (1 Jam Sebelum)",
     "Dikirim: Setiap jam :00 WIB via cron otomatis",
     [
         "• Pesan: '1 jam lagi kamu ada jadwal nih!'",
         "• Detail cabang, jam, layanan, kapster",
         "• Pesan: 'Brangkat sekarang biar santai ya!'",
         "• Customer bisa balas: jadi/telat/cancel/reschedule",
         "• AI bot merespons balasan secara otomatis",
         "• Aturan terlambat maks. 10-15 menit disertakan",
     ]),
]

for i, (color, title, timing, points) in enumerate(notifs):
    x = 0.5 + i * 4.24
    y = 1.65
    w = 4.0

    add_rect(slide, x, y, w, 5.55,
             fill_rgb=MID_GRAY,
             line_rgb=RGBColor(0x66, 0x66, 0x66), line_width=0.5)
    add_rect(slide, x, y, w, 0.52, fill_rgb=color if color != YELLOW else RGBColor(0xB8, 0x86, 0x00))
    add_text(slide, title, x + 0.12, y + 0.06, w - 0.2, 0.42,
             font_size=13, bold=True, color=WHITE if color != YELLOW else BLACK)
    add_text(slide, timing, x + 0.12, y + 0.6, w - 0.2, 0.38,
             font_size=11, italic=True, color=RGBColor(0xAA, 0xEE, 0xAA))

    for j, pt in enumerate(points):
        add_text(slide, pt, x + 0.15, y + 1.08 + j * 0.7, w - 0.25, 0.65,
                 font_size=12, color=WHITE)

# =============================================================================
# SLIDE 5 — TUGAS ADMIN — DASHBOARD
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Dashboard Admin", "Pusat kendali seluruh reservasi semua cabang")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

# Left column — fitur dashboard
left_items = [
    ("📋", "Lihat Semua Booking", "Tampilan tabel semua reservasi hari ini & mendatang"),
    ("🔍", "Filter & Cari", "Filter per cabang, kapster, tanggal, atau status booking"),
    ("✏️", "Edit Booking", "Ubah layanan, tanggal, jam, kapster, atau catatan"),
    ("❌", "Cancel Booking", "Batalkan booking dengan satu klik"),
    ("✅", "Update Status", "Ubah status: pending → confirmed → done → cancel"),
]

add_rect(slide, 0.5, 1.65, 6.0, 5.55,
         fill_rgb=MID_GRAY,
         line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.5)
add_rect(slide, 0.5, 1.65, 6.0, 0.45, fill_rgb=RED)
add_text(slide, "🖥️  Fitur Dashboard Admin", 0.65, 1.68, 5.7, 0.4,
         font_size=15, bold=True, color=WHITE)

for i, (icon, title, desc) in enumerate(left_items):
    iy = 2.25 + i * 0.97
    add_rect(slide, 0.6, iy, 5.75, 0.85,
             fill_rgb=RGBColor(0x38, 0x38, 0x38),
             line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.3)
    add_text(slide, icon, 0.7, iy + 0.05, 0.5, 0.75, font_size=20, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.25, iy + 0.04, 4.9, 0.35, font_size=13, bold=True, color=WHITE)
    add_text(slide, desc, 1.25, iy + 0.42, 4.9, 0.38, font_size=11,
             color=RGBColor(0xBB, 0xBB, 0xBB))

# Right column — status alur
add_rect(slide, 6.9, 1.65, 6.1, 5.55,
         fill_rgb=MID_GRAY,
         line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.5)
add_rect(slide, 6.9, 1.65, 6.1, 0.45, fill_rgb=DARK_RED)
add_text(slide, "🔄  Alur Status Booking", 7.05, 1.68, 5.8, 0.4,
         font_size=15, bold=True, color=WHITE)

statuses = [
    (RGBColor(0x88, 0x88, 0x00), "PENDING",
     "Booking baru masuk, belum dikonfirmasi manual admin"),
    (GREEN,                       "CONFIRMED",
     "Booking dikonfirmasi, kapster disiapkan, WA terkirim"),
    (RGBColor(0x00, 0x70, 0xC0), "DONE",
     "Layanan selesai dilakukan di barbershop"),
    (RED,                         "CANCELLED",
     "Dibatalkan oleh customer atau admin (slot terbuka kembali)"),
    (RGBColor(0x80, 0x00, 0x80), "NO_SHOW",
     "Customer tidak hadir tanpa konfirmasi pembatalan"),
]

for i, (color, status, desc) in enumerate(statuses):
    sy = 2.25 + i * 0.92
    add_rect(slide, 7.05, sy, 1.4, 0.72, fill_rgb=color)
    add_text(slide, status, 7.05, sy, 1.4, 0.72,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, 8.55, sy + 0.1, 4.3, 0.55,
             font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))
    if i < len(statuses) - 1:
        add_text(slide, "↓", 7.6, sy + 0.68, 0.4, 0.3,
                 font_size=11, color=YELLOW, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 6 — TUGAS ADMIN PER CABANG (HARIAN)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Tugas Admin per Cabang", "Rutinitas harian admin di masing-masing cabang RedBox")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

# 5 cabang
branches = [
    ("Bypass\n(Pusat)", "0818-202-569", "10.00–22.00", RED),
    ("Samadikun", "0818-202-589", "10.00–21.00", DARK_RED),
    ("CSB Mall", "0818-202-889", "10.00–21.00", RGBColor(0x8B, 0x45, 0x13)),
    ("Sumber", "0818-202-599", "10.00–21.00", RGBColor(0x00, 0x56, 0x8B)),
    ("Tegal", "0818-268-883", "10.00–21.00", RGBColor(0x2E, 0x6B, 0x2E)),
]

bw = 2.35
bx0 = 0.5
for i, (name, wa, hours, color) in enumerate(branches):
    bx = bx0 + i * (bw + 0.15)
    add_rect(slide, bx, 1.65, bw, 0.7, fill_rgb=color)
    add_text(slide, name, bx, 1.65, bw, 0.7,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, 2.35, bw, 0.35, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_text(slide, f"📞 {wa}  |  🕐 {hours}", bx, 2.35, bw, 0.35,
             font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Tugas harian — shared between branches
add_rect(slide, 0.5, 2.85, 12.5, 4.35,
         fill_rgb=MID_GRAY,
         line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.5)
add_rect(slide, 0.5, 2.85, 12.5, 0.45, fill_rgb=RGBColor(0x33, 0x33, 0x33))
add_text(slide, "📌  Tugas Harian Admin (Semua Cabang)", 0.65, 2.87, 12, 0.4,
         font_size=14, bold=True, color=WHITE)

tasks = [
    ("🌅  Buka Hari",
     ["Cek jadwal booking masuk via WA notifikasi",
      "Pastikan kapster hadir sesuai jadwal",
      "Konfirmasi booking pending jika ada"]),
    ("📲  Saat Booking Masuk",
     ["Terima notif WA: nama, layanan, jam, cabang",
      "Cek ketersediaan slot di dashboard",
      "Update status ke CONFIRMED jika belum otomatis"]),
    ("⚡  Saat Customer Datang",
     ["Cek nama customer di daftar booking",
      "Arahkan ke kapster yang sudah dijadwalkan",
      "Update status ke DONE setelah selesai"]),
    ("🔄  Reschedule / Cancel",
     ["Terima permintaan via WA atau langsung",
      "Update di dashboard admin",
      "Slot terbuka otomatis untuk booking lain"]),
]

cols2 = 2
row2  = 2
tw = 5.85
th = 1.6

for i, (title, pts) in enumerate(tasks):
    col = i % cols2
    row = i // cols2
    tx = 0.65 + col * (tw + 0.4)
    ty = 3.45 + row * (th + 0.2)
    add_rect(slide, tx, ty, tw, th,
             fill_rgb=RGBColor(0x38, 0x38, 0x38),
             line_rgb=RGBColor(0x60, 0x60, 0x60), line_width=0.3)
    add_rect(slide, tx, ty, tw, 0.36, fill_rgb=RED)
    add_text(slide, title, tx + 0.1, ty + 0.04, tw - 0.15, 0.32,
             font_size=13, bold=True, color=WHITE)
    for j, pt in enumerate(pts):
        add_text(slide, f"• {pt}", tx + 0.12, ty + 0.44 + j * 0.37, tw - 0.2, 0.36,
                 font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

# =============================================================================
# SLIDE 7 — FLOW TEKNIS RESERVASI (ARSITEKTUR)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Flow Teknis Reservasi", "Arsitektur sistem end-to-end dari form ke database ke WA")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

# Main flow diagram — horizontal
nodes = [
    (RGBColor(0x00, 0x60, 0x00), "🌐\nWebsite", "booking.html\nForm submit"),
    (RGBColor(0x00, 0x50, 0x80), "⚙️\nAPI Server", "POST /api/bookings\nValidasi + overlap check"),
    (RGBColor(0x60, 0x00, 0x80), "🗄️\nSupabase DB", "Insert bookings\nUpsert customers"),
    (RED,                         "📱\nFonnte WA", "WA konfirmasi\n+ Admin notif"),
    (RGBColor(0x80, 0x50, 0x00), "🏪\nMoka POS", "Bridge ke schedules\nSync jadwal kapster"),
]

nw   = 2.1
nh   = 2.0
gap2 = 0.28
start = 0.5

for i, (color, title, desc) in enumerate(nodes):
    nx = start + i * (nw + gap2)
    ny = 1.75
    add_rect(slide, nx, ny, nw, nh, fill_rgb=color,
             line_rgb=WHITE, line_width=0.5)
    add_text(slide, title, nx, ny + 0.15, nw, 0.85,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, nx + 0.1, ny + 1.05, nw - 0.2, 0.85,
             font_size=11, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        ax = nx + nw + 0.03
        add_text(slide, "→", ax, ny + 0.7, gap2 + 0.12, 0.5,
                 font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Bottom detail boxes
details = [
    (0.5, "📋  Validasi Booking",
     ["✓ Semua field wajib diisi", "✓ Format WA valid (digit only)", "✓ Tidak ada overlap jadwal kapster", "✓ Status langsung CONFIRMED"]),
    (3.78, "🔒  Anti Double Booking",
     ["✓ Cek overlap barber × waktu", "✓ Legacy bookings juga dicek", "✓ Moka open bills diblock", "✓ Stale bills expire otomatis"]),
    (7.06, "📨  WA Gateway (Fonnte)",
     ["✓ Normalisasi nomor 8xxx → 628xxx", "✓ Token device 0818-202-569", "✓ Kirim konfirmasi & admin notif", "✓ Sama dipakai AI bot & reminder"]),
    (10.34, "🏪  Moka POS Sync",
     ["✓ Bridge booking ke schedules", "✓ Tandai slot kapster busy", "✓ GoShow blocks expire 1-4 jam", "✓ Cron expire setiap 15 menit"]),
]

for dx, dtitle, dpts in details:
    add_rect(slide, dx, 4.1, 2.8, 3.1,
             fill_rgb=MID_GRAY,
             line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.3)
    add_rect(slide, dx, 4.1, 2.8, 0.38, fill_rgb=DARK_RED)
    add_text(slide, dtitle, dx + 0.1, 4.12, 2.6, 0.35,
             font_size=11, bold=True, color=WHITE)
    for j, pt in enumerate(dpts):
        add_text(slide, pt, dx + 0.1, 4.58 + j * 0.6, 2.65, 0.56,
                 font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

# =============================================================================
# SLIDE 8 — SISTEM REMINDER & CRON OTOMATIS
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Sistem Reminder & Cron Otomatis", "4 cron jobs yang berjalan otomatis di Vercel")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

crons = [
    (GREEN,                        "⏰  Remind-Soon\n(1 Jam Sebelum)",
     "Jadwal:  Setiap jam :00 WIB\nEndpoint: /api/cron/remind-soon",
     ["Mencari booking dengan jam = jam berikutnya",
      "Kirim WA: '1 jam lagi ada jadwal kamu!'",
      "Include: cabang, jam, layanan, kapster",
      "Customer bisa balas: jadi / telat / cancel",
      "AI bot merespons balasan secara otomatis"]),
    (ORANGE,                       "📅  Reminder H-1\n(Sehari Sebelum)",
     "Jadwal:  Setiap hari jam 10:00 WIB\nEndpoint: /api/cron/reminders",
     ["Mencari semua booking besok (semua cabang)",
      "Kirim WA reminder dengan detail jadwal",
      "Include: tanggal, jam, layanan, kapster, cabang",
      "Tidak dikirim jika booking sudah cancelled",
      "Berlaku otomatis untuk 5 cabang sekaligus"]),
    (RGBColor(0x00, 0x70, 0xC0),   "🗑️  Expire Stale Bills\n(Moka GoShow)",
     "Jadwal:  Setiap 15 menit\nEndpoint: /api/cron/expire-stale-bills",
     ["Expire null-barber blocks jika end_time+1j lewat",
      "Expire null-barber blocks jika created_at+2j lewat",
      "Expire per-barber blocks jika end_time+4j lewat",
      "Mencegah GoShow Moka blokir slot lintas hari",
      "Configurable via env: MOKA_OPENBILL_STALE_HOURS"]),
    (RGBColor(0x60, 0x00, 0x80),   "🎂  Birthday Reminder",
     "Jadwal:  Setiap hari jam 08:00 WIB\nEndpoint: /api/cron/birthday",
     ["Cek customer dengan ulang tahun hari ini",
      "Kirim WA ucapan ulang tahun personal",
      "Optionally: sertakan promo ulang tahun",
      "Data ulang tahun dari tabel customers",
      "Fire-and-forget, non-blocking"]),
]

cw2 = 3.0
gap3 = 0.2
start2 = 0.48

for i, (color, title, schedule, pts) in enumerate(crons):
    cx2 = start2 + i * (cw2 + gap3)
    add_rect(slide, cx2, 1.65, cw2, 5.55,
             fill_rgb=MID_GRAY,
             line_rgb=RGBColor(0x55, 0x55, 0x55), line_width=0.5)
    add_rect(slide, cx2, 1.65, cw2, 0.7, fill_rgb=color)
    add_text(slide, title, cx2 + 0.1, 1.67, cw2 - 0.15, 0.65,
             font_size=13, bold=True, color=WHITE if color != YELLOW else BLACK,
             align=PP_ALIGN.CENTER)
    add_rect(slide, cx2, 2.35, cw2, 0.55, fill_rgb=RGBColor(0x1A, 0x1A, 0x1A))
    add_text(slide, schedule, cx2 + 0.1, 2.37, cw2 - 0.15, 0.52,
             font_size=10, italic=True, color=RGBColor(0xAA, 0xDD, 0xFF))
    for j, pt in enumerate(pts):
        add_text(slide, f"• {pt}", cx2 + 0.12, 3.02 + j * 0.7, cw2 - 0.2, 0.66,
                 font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

# =============================================================================
# SLIDE 9 — LIFECYCLE & STATUS BOOKING LENGKAP
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, DARK_GRAY)
top_bar(slide, "Lifecycle & Status Booking", "Alur perubahan status dari booking dibuat hingga selesai")
add_rect(slide, 0, 0, 0.35, 7.5, fill_rgb=RED)

# Timeline flow
timeline = [
    (GREEN,                      "CONFIRMED",   "Booking dibuat\nvia website",         "Real-time\nsaat submit"),
    (RGBColor(0x00, 0x70, 0xC0), "WA KONFIRMASI", "Terkirim ke\ncustomer",             "Detik setelah\nbooking"),
    (ORANGE,                     "REMINDER H-1", "WA reminder\nsehari sebelum",         "10:00 WIB\nH-1"),
    (YELLOW,                     "REMIND-SOON",  "WA '1 jam lagi'\nke customer",        "1 jam\nsebelum"),
    (RGBColor(0x00, 0x56, 0x8B), "CUSTOMER\nDATANG",  "Kapster siap\nmelayani",         "Sesuai\njam booking"),
    (GREEN,                      "DONE",         "Layanan selesai\nAdmin update status", "Saat layanan\nselesai"),
]

tw2 = 1.75
tgap = 0.25
tx0 = 0.55

for i, (color, status, desc, timing) in enumerate(timeline):
    tx2 = tx0 + i * (tw2 + tgap)
    # Main status box
    add_rect(slide, tx2, 1.8, tw2, 1.0, fill_rgb=color,
             line_rgb=WHITE, line_width=0.5)
    add_text(slide, status, tx2, 1.8, tw2, 1.0,
             font_size=13, bold=True, color=WHITE if color != YELLOW else BLACK,
             align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(timeline) - 1:
        ax2 = tx2 + tw2 + 0.03
        add_text(slide, "→", ax2, 2.1, tgap + 0.1, 0.4,
                 font_size=16, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    # Description below
    add_text(slide, desc, tx2, 2.9, tw2, 0.6,
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # Timing badge
    add_rect(slide, tx2, 3.58, tw2, 0.55, fill_rgb=RGBColor(0x22, 0x22, 0x22))
    add_text(slide, timing, tx2, 3.58, tw2, 0.55,
             font_size=10, italic=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Alternative endings
add_rect(slide, 0.55, 4.4, 12.5, 0.04, fill_rgb=RGBColor(0x55, 0x55, 0x55))
add_text(slide, "Alternatif Status:", 0.55, 4.55, 3, 0.35,
         font_size=14, bold=True, color=WHITE)

alt_statuses = [
    (RED,                         "CANCELLED",
     "Customer atau admin batalkan sebelum layanan · Slot terbuka kembali otomatis"),
    (RGBColor(0x80, 0x00, 0x80), "NO_SHOW",
     "Customer tidak hadir tanpa konfirmasi · Admin mark secara manual"),
]

for i, (color, st, desc) in enumerate(alt_statuses):
    ax3 = 0.55 + i * 6.3
    add_rect(slide, ax3, 5.0, 5.9, 0.85,
             fill_rgb=RGBColor(0x35, 0x35, 0x35),
             line_rgb=color, line_width=1.5)
    add_rect(slide, ax3, 5.0, 1.6, 0.85, fill_rgb=color)
    add_text(slide, st, ax3, 5.0, 1.6, 0.85,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, ax3 + 1.7, 5.1, 4.1, 0.65,
             font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

# WA bot response note
add_rect(slide, 0.55, 6.05, 12.5, 1.15,
         fill_rgb=RGBColor(0x00, 0x35, 0x55),
         line_rgb=RGBColor(0x00, 0x88, 0xCC), line_width=1)
add_text(slide, "💬  AI Bot WhatsApp (Reddy)", 0.75, 6.1, 5, 0.38,
         font_size=14, bold=True, color=WHITE)
add_text(slide,
         "Merespons balasan customer terhadap reminder secara otomatis:\n"
         "• Konfirmasi hadir → '\"Sip, ditunggu kak!\" + aturan keterlambatan'\n"
         "• Telat/macet → empati + aturan terlambat maks 10-15 menit'\n"
         "• Cancel → empati + tawarkan reschedule via link booking",
         0.75, 6.52, 12, 0.68, font_size=11, color=RGBColor(0xCC, 0xEE, 0xFF))

# =============================================================================
# SLIDE 10 — PENUTUP / RINGKASAN
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, BLACK)
add_rect(slide, 0, 0, 0.5, 7.5, fill_rgb=RED)
add_rect(slide, 0.5, 3.65, 12.83, 0.06, fill_rgb=RED)

add_text(slide, "RINGKASAN SISTEM", 1.0, 0.5, 11.8, 0.9,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "RedBox Barbershop Reservation Flow", 1.0, 1.35, 11.8, 0.5,
         font_size=18, color=RED, align=PP_ALIGN.LEFT)

summary_items = [
    ("Customer", [
        "Booking via website — 6 langkah mudah",
        "Slot real-time tanpa double booking",
        "Terima 3x WA otomatis (konfirmasi, H-1, 1jam)",
    ]),
    ("Admin", [
        "Terima notif WA setiap booking baru",
        "Kelola di dashboard: edit/cancel/update status",
        "Koordinasi kapster per cabang",
    ]),
    ("Sistem", [
        "4 cron otomatis: reminder, expire, birthday",
        "Integrasi Moka POS untuk slot kapster",
        "AI bot WA Reddy merespons 24/7",
    ]),
]

for i, (title, pts) in enumerate(summary_items):
    sx = 1.0 + i * 4.1
    add_rect(slide, sx, 2.0, 3.8, 1.45, fill_rgb=RED)
    add_text(slide, title, sx, 2.0, 3.8, 1.45,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for j, pt in enumerate(pts):
        add_text(slide, f"→  {pt}", sx, 3.85 + j * 0.52, 3.8, 0.5,
                 font_size=13, color=WHITE)

add_rect(slide, 0.5, 6.3, 12.83, 0.05, fill_rgb=RGBColor(0x44, 0x44, 0x44))
add_text(slide,
         "redboxbarbershop.com  ·  WA Bypass: 0818-202-569  ·  5 Cabang di Cirebon & Tegal",
         0.5, 6.45, 12.8, 0.4,
         font_size=13, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# =============================================================================
# SAVE
# =============================================================================
out = r"c:\Users\Win11\Downloads\Documents\Digital Market\Website RedBox\RedBox_Reservasi_Flow.pptx"
prs.save(out)
print(f"✅  Presentasi berhasil disimpan: {out}")
print(f"    Total slides: {len(prs.slides)}")
